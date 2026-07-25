"""
Graphico Pro - Business Logic Services
Core service layer for AI operations, file management, and project handling.
"""

import json
import logging
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple

from config.settings import (
    OPENAI_API_KEY,
    ANTHROPIC_API_KEY,
    GOOGLE_AI_API_KEY,
    STABILITY_API_KEY,
    REPLICATE_API_KEY,
    DEEPSEEK_API_KEY,
    GITHUB_MODELS_API_KEY,
    MAX_PROJECTS_PER_USER,
    MAX_FILES_PER_PROJECT,
    MAX_CHAT_HISTORY_LENGTH,
    MAX_UPLOAD_SIZE_BYTES,
    ALLOWED_UPLOAD_EXTENSIONS,
    UPLOAD_DIRECTORY,
    EXPORT_DIRECTORY,
    TEMP_DIRECTORY,
)
from config.constants import (
    AIProvider,
    AIModel,
    ProjectStatus,
    FileCategory,
    HttpStatus,
    ApiMessage,
    ErrorCode,
    SystemLimits,
    DefaultPrompts,
    Pagination,
    MIME_TYPES,
)
from backend.utils import (
    generate_id,
    get_utc_now,
    get_timestamp,
    read_json_file,
    write_json_file,
    success_response,
    error_response,
    validate_file_extension,
    validate_file_size,
    sanitize_filename,
    get_file_category,
    get_file_extension,
    get_mime_type,
    humanize_bytes,
    ensure_directory,
    clean_temp_files,
    safe_path_join,
    log_execution_time,
    logger as utils_logger,
)

logger = logging.getLogger(__name__)


# ============================================================================
# LAZY IMPORT HELPERS
# ============================================================================

def _import_openai():
    try:
        import openai
        return openai
    except ImportError:
        logger.warning("OpenAI package not installed")
        return None


def _import_anthropic():
    try:
        import anthropic
        return anthropic
    except ImportError:
        logger.warning("Anthropic package not installed")
        return None


def _import_google_ai():
    try:
        import google.generativeai as genai
        return genai
    except ImportError:
        logger.warning("Google AI package not installed")
        return None


def _import_pypdf():
    try:
        import PyPDF2
        return PyPDF2
    except ImportError:
        logger.debug("PyPDF2 not installed")
        return None


def _import_pdfplumber():
    try:
        import pdfplumber
        return pdfplumber
    except ImportError:
        logger.debug("pdfplumber not installed")
        return None


# ============================================================================
# AI SERVICE
# ============================================================================

class AIService:
    """
    Manages interactions with various AI providers.
    Provides unified interface for chat, image generation, and document processing.
    """

    def __init__(self):
        self.providers = {}
        self._init_providers()
        self.available_models = self._get_available_models()
        logger.info(f"AI Service initialized with {len(self.available_models)} available models")

    def _init_providers(self) -> None:
        openai_lib = _import_openai()

        # GitHub Models (free GPT-4o via Azure — listed first so it's default)
        if GITHUB_MODELS_API_KEY:
            if openai_lib:
                try:
                    self.providers[AIProvider.GITHUB_MODELS.value] = openai_lib.OpenAI(
                        api_key=GITHUB_MODELS_API_KEY,
                        base_url="https://models.inference.ai.azure.com"
                    )
                    logger.info("GitHub Models client initialized successfully")
                except Exception as e:
                    logger.error(f"Failed to initialize GitHub Models client: {e}")
            else:
                logger.warning("OpenAI package not available for GitHub Models - skipping")
        else:
            logger.info("GitHub Models API key not configured - skipping")

        # OpenAI
        if OPENAI_API_KEY:
            if openai_lib:
                try:
                    self.providers[AIProvider.OPENAI.value] = openai_lib.OpenAI(api_key=OPENAI_API_KEY)
                    logger.info("OpenAI client initialized successfully")
                except Exception as e:
                    logger.error(f"Failed to initialize OpenAI client: {e}")
            else:
                logger.warning("OpenAI package not available - skipping")
        else:
            logger.info("OpenAI API key not configured - skipping")

        # Anthropic
        if ANTHROPIC_API_KEY:
            anthropic = _import_anthropic()
            if anthropic:
                try:
                    self.providers[AIProvider.ANTHROPIC.value] = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
                    logger.info("Anthropic client initialized successfully")
                except Exception as e:
                    logger.error(f"Failed to initialize Anthropic client: {e}")
            else:
                logger.warning("Anthropic package not available - skipping")
        else:
            logger.info("Anthropic API key not configured - skipping")

        # Google AI
        if GOOGLE_AI_API_KEY:
            genai = _import_google_ai()
            if genai:
                try:
                    genai.configure(api_key=GOOGLE_AI_API_KEY)
                    self.providers[AIProvider.GOOGLE_AI.value] = genai
                    logger.info("Google AI client initialized successfully")
                except Exception as e:
                    logger.error(f"Failed to initialize Google AI client: {e}")
            else:
                logger.warning("Google AI package not available - skipping")
        else:
            logger.info("Google AI API key not configured - skipping")

        # Stability AI
        if STABILITY_API_KEY:
            self.providers[AIProvider.STABILITY.value] = {
                "api_key": STABILITY_API_KEY,
                "configured": True,
            }
            logger.info("Stability AI configured")

        # Replicate
        if REPLICATE_API_KEY:
            self.providers[AIProvider.REPLICATE.value] = {
                "api_key": REPLICATE_API_KEY,
                "configured": True,
            }
            logger.info("Replicate configured")

        # DeepSeek (OpenAI-compatible API)
        if DEEPSEEK_API_KEY:
            if openai_lib:
                try:
                    self.providers[AIProvider.DEEPSEEK.value] = openai_lib.OpenAI(
                        api_key=DEEPSEEK_API_KEY,
                        base_url="https://api.deepseek.com"
                    )
                    logger.info("DeepSeek client initialized successfully")
                except Exception as e:
                    logger.error(f"Failed to initialize DeepSeek client: {e}")
            else:
                logger.warning("OpenAI package not available for DeepSeek - skipping")
        else:
            logger.info("DeepSeek API key not configured - skipping")

    def _get_available_models(self) -> List[Dict[str, Any]]:
        models = []

        # GitHub Models (listed first — becomes default)
        if AIProvider.GITHUB_MODELS.value in self.providers:
            models.extend([
                {
                    "model": AIModel.GITHUB_GPT4O.value,
                    "provider": AIProvider.GITHUB_MODELS.value,
                    "type": "chat",
                    "display_name": "GPT-4o (GitHub)",
                    "description": "Free GPT-4o via GitHub Marketplace",
                    "max_tokens": 128000,
                },
            ])

        # DeepSeek Models
        if AIProvider.DEEPSEEK.value in self.providers:
            models.extend([
                {
                    "model": AIModel.DEEPSEEK_CHAT.value,
                    "provider": AIProvider.DEEPSEEK.value,
                    "type": "chat",
                    "display_name": "DeepSeek Chat",
                    "description": "Fast general-purpose chat model",
                    "max_tokens": 32000,
                },
                {
                    "model": AIModel.DEEPSEEK_REASONER.value,
                    "provider": AIProvider.DEEPSEEK.value,
                    "type": "chat",
                    "display_name": "DeepSeek Reasoner",
                    "description": "Advanced reasoning model for complex tasks",
                    "max_tokens": 32000,
                },
            ])

        # OpenAI Models
        if AIProvider.OPENAI.value in self.providers:
            models.extend([
                {
                    "model": AIModel.GPT_4O.value,
                    "provider": AIProvider.OPENAI.value,
                    "type": "chat",
                    "display_name": "GPT-4o",
                    "description": "Most capable OpenAI model with vision",
                    "max_tokens": 128000,
                },
                {
                    "model": AIModel.GPT_4O_MINI.value,
                    "provider": AIProvider.OPENAI.value,
                    "type": "chat",
                    "display_name": "GPT-4o Mini",
                    "description": "Fast and efficient OpenAI model",
                    "max_tokens": 128000,
                },
                {
                    "model": AIModel.GPT_4_TURBO.value,
                    "provider": AIProvider.OPENAI.value,
                    "type": "chat",
                    "display_name": "GPT-4 Turbo",
                    "description": "Powerful GPT-4 with extended context",
                    "max_tokens": 128000,
                },
                {
                    "model": AIModel.GPT_35_TURBO.value,
                    "provider": AIProvider.OPENAI.value,
                    "type": "chat",
                    "display_name": "GPT-3.5 Turbo",
                    "description": "Cost-effective chat model",
                    "max_tokens": 16385,
                },
                {
                    "model": AIModel.DALL_E_3.value,
                    "provider": AIProvider.OPENAI.value,
                    "type": "image",
                    "display_name": "DALL-E 3",
                    "description": "Advanced image generation with prompt understanding",
                },
                {
                    "model": AIModel.DALL_E_2.value,
                    "provider": AIProvider.OPENAI.value,
                    "type": "image",
                    "display_name": "DALL-E 2",
                    "description": "Fast image generation",
                },
            ])

        # Anthropic Models
        if AIProvider.ANTHROPIC.value in self.providers:
            models.extend([
                {
                    "model": AIModel.CLAUDE_35_SONNET.value,
                    "provider": AIProvider.ANTHROPIC.value,
                    "type": "chat",
                    "display_name": "Claude 3.5 Sonnet",
                    "description": "Most intelligent Claude model",
                    "max_tokens": 200000,
                },
                {
                    "model": AIModel.CLAUDE_3_OPUS.value,
                    "provider": AIProvider.ANTHROPIC.value,
                    "type": "chat",
                    "display_name": "Claude 3 Opus",
                    "description": "Powerful Claude model for complex tasks",
                    "max_tokens": 200000,
                },
                {
                    "model": AIModel.CLAUDE_3_SONNET.value,
                    "provider": AIProvider.ANTHROPIC.value,
                    "type": "chat",
                    "display_name": "Claude 3 Sonnet",
                    "description": "Balanced Claude model",
                    "max_tokens": 200000,
                },
                {
                    "model": AIModel.CLAUDE_3_HAIKU.value,
                    "provider": AIProvider.ANTHROPIC.value,
                    "type": "chat",
                    "display_name": "Claude 3 Haiku",
                    "description": "Fastest Claude model for quick tasks",
                    "max_tokens": 200000,
                },
            ])

        # Google AI Models
        if AIProvider.GOOGLE_AI.value in self.providers:
            models.extend([
                {
                    "model": AIModel.GEMINI_20_FLASH.value,
                    "provider": AIProvider.GOOGLE_AI.value,
                    "type": "chat",
                    "display_name": "Gemini 2.0 Flash",
                    "description": "Fast and reliable Gemini model",
                    "max_tokens": 1000000,
                },
                {
                    "model": AIModel.GEMINI_25_PRO.value,
                    "provider": AIProvider.GOOGLE_AI.value,
                    "type": "chat",
                    "display_name": "Gemini 2.5 Pro",
                    "description": "Most capable Gemini model",
                    "max_tokens": 2000000,
                },
                {
                    "model": AIModel.GEMINI_25_FLASH.value,
                    "provider": AIProvider.GOOGLE_AI.value,
                    "type": "chat",
                    "display_name": "Gemini 2.5 Flash",
                    "description": "Fast and efficient Gemini model",
                    "max_tokens": 1000000,
                },
            ])

        # Stability AI Models
        if AIProvider.STABILITY.value in self.providers:
            models.extend([
                {
                    "model": AIModel.STABLE_DIFFUSION_XL.value,
                    "provider": AIProvider.STABILITY.value,
                    "type": "image",
                    "display_name": "Stable Diffusion XL",
                    "description": "High-quality open-source image generation",
                },
                {
                    "model": AIModel.STABLE_DIFFUSION_3.value,
                    "provider": AIProvider.STABILITY.value,
                    "type": "image",
                    "display_name": "Stable Diffusion 3",
                    "description": "Latest Stable Diffusion model",
                },
            ])

        # Replicate Models
        if AIProvider.REPLICATE.value in self.providers:
            models.extend([
                {
                    "model": AIModel.SDXL_LIGHTNING.value,
                    "provider": AIProvider.REPLICATE.value,
                    "type": "image",
                    "display_name": "SDXL Lightning",
                    "description": "Ultra-fast image generation on Replicate",
                },
            ])

        return models

    def get_available_models(self) -> List[Dict[str, Any]]:
        return self.available_models.copy()

    def get_model_info(self, model_name: str) -> Optional[Dict[str, Any]]:
        for model in self.available_models:
            if model["model"] == model_name:
                return model.copy()
        return None

    def is_provider_available(self, provider: str) -> bool:
        return provider in self.providers

    @log_execution_time
    def chat(
        self,
        messages: List[Dict[str, str]],
        model: str = None,
        system_prompt: Optional[str] = None,
        max_tokens: int = 2000,
        temperature: float = 0.7,
        **kwargs,
    ) -> Tuple[Optional[Dict[str, Any]], Optional[Tuple]]:
        try:
            if model is None:
                if self.available_models:
                    model = self.available_models[0]["model"]
                else:
                    model = AIModel.GITHUB_GPT4O.value

            model_info = self.get_model_info(model)

            if not model_info:
                return None, error_response(
                    message=ApiMessage.AI_MODEL_UNAVAILABLE,
                    error_code=ErrorCode.AI_MODEL_ERROR,
                    status_code=HttpStatus.BAD_REQUEST,
                    error=f"Model not found or not available: {model}",
                )

            provider = model_info["provider"]

            if not self.is_provider_available(provider):
                return None, error_response(
                    message=ApiMessage.AI_MODEL_UNAVAILABLE,
                    error_code=ErrorCode.AI_PROVIDER_ERROR,
                    status_code=HttpStatus.SERVICE_UNAVAILABLE,
                    error=f"Provider {provider} is not configured",
                )

            temperature = max(0.0, min(1.0, temperature))

            full_messages = []
            if system_prompt:
                full_messages.append({"role": "system", "content": system_prompt})
            full_messages.extend(messages[-SystemLimits.MAX_CONTEXT_WINDOW:])

            if not full_messages:
                return None, error_response(
                    message="No messages provided",
                    error_code=ErrorCode.VALIDATION_MISSING_FIELD,
                    status_code=HttpStatus.BAD_REQUEST,
                )

            # OpenAI-compatible providers (OpenAI, DeepSeek, GitHub Models)
            if provider in (AIProvider.OPENAI.value, AIProvider.DEEPSEEK.value, AIProvider.GITHUB_MODELS.value):
                return self._chat_openai(full_messages, model, max_tokens, temperature, **kwargs)
            elif provider == AIProvider.ANTHROPIC.value:
                return self._chat_anthropic(full_messages, model, max_tokens, temperature, **kwargs)
            elif provider == AIProvider.GOOGLE_AI.value:
                return self._chat_google(full_messages, model, max_tokens, temperature, **kwargs)
            else:
                return None, error_response(
                    message="Provider not supported for chat",
                    error_code=ErrorCode.AI_PROVIDER_ERROR,
                    status_code=HttpStatus.BAD_REQUEST,
                )

        except Exception as e:
            logger.error(f"Chat generation error: {e}", exc_info=True)
            return None, error_response(
                message=ApiMessage.AI_GENERATION_FAILED,
                error_code=ErrorCode.AI_GENERATION_ERROR,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def _chat_openai(
        self,
        messages: List[Dict[str, str]],
        model: str,
        max_tokens: int,
        temperature: float,
        **kwargs,
    ) -> Tuple[Dict[str, Any], None]:
        """Handle chat completion with OpenAI and OpenAI-compatible APIs."""
        model_info = self.get_model_info(model)
        provider = model_info["provider"] if model_info else AIProvider.OPENAI.value
        client = self.providers[provider]

        response = client.chat.completions.create(
            model=model,
            messages=messages,
            max_tokens=max_tokens,
            temperature=temperature,
            **kwargs,
        )

        result = {
            "model": model,
            "provider": provider,
            "content": response.choices[0].message.content,
            "role": response.choices[0].message.role,
            "usage": {
                "prompt_tokens": response.usage.prompt_tokens if response.usage else 0,
                "completion_tokens": response.usage.completion_tokens if response.usage else 0,
                "total_tokens": response.usage.total_tokens if response.usage else 0,
            },
            "finish_reason": response.choices[0].finish_reason,
        }

        logger.debug(f"{provider} chat completed: {result['usage']['total_tokens']} tokens used")
        return result, None

    def _chat_anthropic(
        self,
        messages: List[Dict[str, str]],
        model: str,
        max_tokens: int,
        temperature: float,
        **kwargs,
    ) -> Tuple[Dict[str, Any], None]:
        """Handle chat completion with Anthropic Claude."""
        client = self.providers[AIProvider.ANTHROPIC.value]
        system_prompt = None
        chat_messages = []
        for msg in messages:
            if msg["role"] == "system":
                system_prompt = msg["content"]
            else:
                chat_messages.append(msg)

        api_params = {
            "model": model,
            "messages": chat_messages,
            "max_tokens": max_tokens,
            "temperature": temperature,
            **kwargs,
        }
        if system_prompt:
            api_params["system"] = system_prompt

        response = client.messages.create(**api_params)
        result = {
            "model": model,
            "provider": AIProvider.ANTHROPIC.value,
            "content": response.content[0].text,
            "role": response.role,
            "usage": {
                "input_tokens": response.usage.input_tokens,
                "output_tokens": response.usage.output_tokens,
                "total_tokens": response.usage.input_tokens + response.usage.output_tokens,
            },
            "finish_reason": response.stop_reason,
        }
        logger.debug(f"Anthropic chat completed: {result['usage']['total_tokens']} tokens used")
        return result, None

    def _chat_google(
        self,
        messages: List[Dict[str, str]],
        model: str,
        max_tokens: int,
        temperature: float,
        **kwargs,
    ) -> Tuple[Dict[str, Any], None]:
        """Handle chat completion with Google AI Gemini."""
        genai = self.providers[AIProvider.GOOGLE_AI.value]

        system_prompt = None
        for msg in messages:
            if msg["role"] == "system":
                system_prompt = msg["content"]
                break

        history = []
        for msg in messages[:-1]:
            if msg["role"] == "user":
                history.append({"role": "user", "parts": [msg["content"]]})
            elif msg["role"] == "assistant":
                history.append({"role": "model", "parts": [msg["content"]]})

        last_message = messages[-1]["content"] if messages else "Hello"

        generation_config = {
            "max_output_tokens": max_tokens,
            "temperature": temperature,
        }

        model_instance = genai.GenerativeModel(
            model_name=model,
            system_instruction=system_prompt,
            generation_config=generation_config,
        )

        chat_session = model_instance.start_chat(history=history)
        response = chat_session.send_message(last_message)

        result = {
            "model": model,
            "provider": AIProvider.GOOGLE_AI.value,
            "content": response.text,
            "role": "assistant",
            "usage": {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0},
            "finish_reason": "stop",
        }

        try:
            if hasattr(response, 'usage_metadata'):
                result["usage"]["prompt_tokens"] = getattr(response.usage_metadata, 'prompt_token_count', 0)
                result["usage"]["completion_tokens"] = getattr(response.usage_metadata, 'candidates_token_count', 0)
                result["usage"]["total_tokens"] = result["usage"]["prompt_tokens"] + result["usage"]["completion_tokens"]
        except Exception:
            pass

        logger.debug("Google AI chat completed")
        return result, None

    @log_execution_time
    def generate_image(
        self,
        prompt: str,
        model: str = None,
        size: str = "1024x1024",
        quality: str = "standard",
        n: int = 1,
        **kwargs,
    ) -> Tuple[Optional[Dict[str, Any]], Optional[Tuple]]:
        try:
            if model is None:
                model = AIModel.DALL_E_3.value

            model_info = self.get_model_info(model)
            if not model_info:
                return None, error_response(
                    message=ApiMessage.AI_MODEL_UNAVAILABLE,
                    error_code=ErrorCode.AI_MODEL_ERROR,
                    status_code=HttpStatus.BAD_REQUEST,
                )

            provider = model_info["provider"]
            if not self.is_provider_available(provider):
                return None, error_response(
                    message=ApiMessage.AI_MODEL_UNAVAILABLE,
                    error_code=ErrorCode.AI_PROVIDER_ERROR,
                    status_code=HttpStatus.SERVICE_UNAVAILABLE,
                )

            if n > SystemLimits.MAX_IMAGE_GENERATION_BATCH:
                n = SystemLimits.MAX_IMAGE_GENERATION_BATCH

            if provider == AIProvider.OPENAI.value:
                client = self.providers[AIProvider.OPENAI.value]
                if model == AIModel.DALL_E_3.value:
                    n = 1

                response = client.images.generate(
                    model=model, prompt=prompt, size=size, quality=quality, n=n, **kwargs,
                )

                images = []
                for img in response.data:
                    images.append({
                        "url": img.url,
                        "revised_prompt": getattr(img, 'revised_prompt', prompt),
                    })

                result = {
                    "model": model,
                    "provider": AIProvider.OPENAI.value,
                    "images": images,
                    "prompt": prompt,
                }
                logger.info(f"Image generated with {model}: {len(images)} image(s)")
                return result, None

            elif provider in (AIProvider.STABILITY.value, AIProvider.REPLICATE.value):
                return None, error_response(
                    message=f"{provider} image generation available via REST API",
                    error_code=ErrorCode.NOT_IMPLEMENTED,
                    status_code=HttpStatus.NOT_IMPLEMENTED,
                )
            else:
                return None, error_response(
                    message="Provider not supported for image generation",
                    error_code=ErrorCode.AI_PROVIDER_ERROR,
                    status_code=HttpStatus.BAD_REQUEST,
                )

        except Exception as e:
            logger.error(f"Image generation error: {e}", exc_info=True)
            return None, error_response(
                message=ApiMessage.AI_GENERATION_FAILED,
                error_code=ErrorCode.AI_GENERATION_ERROR,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )


# ============================================================================
# PROJECT SERVICE
# ============================================================================

class ProjectService:
    """
    Manages user projects including creation, retrieval, updates, and deletion.
    """

    def __init__(self):
        from config.settings import JSON_PROJECTS_FILE
        self.projects_file = JSON_PROJECTS_FILE
        self._projects: Dict[str, Dict[str, Any]] = {}
        self._load_projects()

    def _load_projects(self) -> None:
        try:
            ensure_directory(self.projects_file.parent)
            self._projects = read_json_file(self.projects_file, {})
            logger.info(f"Loaded {len(self._projects)} projects from storage")
        except Exception as e:
            logger.error(f"Failed to load projects: {e}")
            self._projects = {}

    def _save_projects(self) -> bool:
        ensure_directory(self.projects_file.parent)
        return write_json_file(self.projects_file, self._projects)

    def create_project(
        self, user_id: str, name: str, description: str = "", project_type: str = "general",
    ) -> Tuple[Optional[Dict[str, Any]], Optional[Tuple]]:
        try:
            if not name or not name.strip():
                return None, error_response(
                    message="Project name is required",
                    error_code=ErrorCode.VALIDATION_MISSING_FIELD,
                    status_code=HttpStatus.BAD_REQUEST,
                )
            user_projects = self.get_user_projects(user_id)
            if len(user_projects) >= MAX_PROJECTS_PER_USER:
                return None, error_response(
                    message=ApiMessage.PROJECT_LIMIT_REACHED,
                    error_code=ErrorCode.RESOURCE_LIMIT_REACHED,
                    status_code=HttpStatus.FORBIDDEN,
                )
            project_id = generate_id("proj")
            project = {
                "project_id": project_id,
                "user_id": user_id,
                "name": name.strip(),
                "description": description.strip() if description else "",
                "project_type": project_type,
                "status": ProjectStatus.ACTIVE.value,
                "files": [],
                "history": [],
                "metadata": {},
                "settings": {},
                "created_at": get_timestamp(),
                "updated_at": get_timestamp(),
            }
            self._projects[project_id] = project
            self._save_projects()
            logger.info(f"Project created: {project_id} by user {user_id}")
            return project.copy(), None
        except Exception as e:
            logger.error(f"Failed to create project: {e}", exc_info=True)
            return None, error_response(
                message=ApiMessage.INTERNAL_ERROR,
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def get_project(self, project_id: str) -> Optional[Dict[str, Any]]:
        project = self._projects.get(project_id)
        return project.copy() if project else None

    def get_user_projects(
        self, user_id: str, status: Optional[str] = None,
        page: int = Pagination.DEFAULT_PAGE, per_page: int = Pagination.DEFAULT_PER_PAGE,
    ) -> List[Dict[str, Any]]:
        projects = []
        for project in self._projects.values():
            if project.get("user_id") == user_id:
                if status is None or project.get("status") == status:
                    projects.append(project.copy())
        projects.sort(key=lambda x: x.get("updated_at", ""), reverse=True)
        start = (page - 1) * per_page
        return projects[start:start + per_page]

    def get_user_projects_count(self, user_id: str, status: Optional[str] = None) -> int:
        count = 0
        for project in self._projects.values():
            if project.get("user_id") == user_id:
                if status is None or project.get("status") == status:
                    count += 1
        return count

    def update_project(
        self, project_id: str, updates: Dict[str, Any],
    ) -> Tuple[Optional[Dict[str, Any]], Optional[Tuple]]:
        try:
            project = self._projects.get(project_id)
            if not project:
                return None, error_response(
                    message=ApiMessage.PROJECT_NOT_FOUND,
                    error_code=ErrorCode.RESOURCE_NOT_FOUND,
                    status_code=HttpStatus.NOT_FOUND,
                )
            for field in ["project_id", "user_id", "created_at"]:
                updates.pop(field, None)
            project.update(updates)
            project["updated_at"] = get_timestamp()
            self._projects[project_id] = project
            self._save_projects()
            logger.info(f"Project updated: {project_id}")
            return project.copy(), None
        except Exception as e:
            logger.error(f"Failed to update project: {e}", exc_info=True)
            return None, error_response(
                message=ApiMessage.INTERNAL_ERROR,
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def delete_project(self, project_id: str) -> Tuple[bool, Optional[Tuple]]:
        try:
            if project_id not in self._projects:
                return False, error_response(
                    message=ApiMessage.PROJECT_NOT_FOUND,
                    error_code=ErrorCode.RESOURCE_NOT_FOUND,
                    status_code=HttpStatus.NOT_FOUND,
                )
            project = self._projects[project_id]
            for file_info in project.get("files", []):
                file_path_str = file_info.get("path", "")
                if file_path_str:
                    file_path = Path(file_path_str)
                    if file_path.exists():
                        try:
                            file_path.unlink()
                            logger.debug(f"Deleted associated file: {file_path}")
                        except OSError as e:
                            logger.warning(f"Could not delete file {file_path}: {e}")
            del self._projects[project_id]
            self._save_projects()
            logger.info(f"Project deleted: {project_id}")
            return True, None
        except Exception as e:
            logger.error(f"Failed to delete project: {e}", exc_info=True)
            return False, error_response(
                message=ApiMessage.INTERNAL_ERROR,
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def archive_project(self, project_id: str) -> Tuple[bool, Optional[Tuple]]:
        result, _ = self.update_project(project_id, {"status": ProjectStatus.ARCHIVED.value})
        return (result is not None), None

    def add_file_to_project(
        self, project_id: str, file_info: Dict[str, Any],
    ) -> Tuple[bool, Optional[Tuple]]:
        project = self._projects.get(project_id)
        if not project:
            return False, error_response(
                message=ApiMessage.PROJECT_NOT_FOUND,
                error_code=ErrorCode.RESOURCE_NOT_FOUND,
                status_code=HttpStatus.NOT_FOUND,
            )
        if len(project.get("files", [])) >= MAX_FILES_PER_PROJECT:
            return False, error_response(
                message="Maximum files per project reached",
                error_code=ErrorCode.RESOURCE_LIMIT_REACHED,
                status_code=HttpStatus.FORBIDDEN,
            )
        project.setdefault("files", []).append(file_info)
        project["updated_at"] = get_timestamp()
        self._projects[project_id] = project
        self._save_projects()
        return True, None

    def remove_file_from_project(
        self, project_id: str, file_id: str,
    ) -> Tuple[bool, Optional[Tuple]]:
        project = self._projects.get(project_id)
        if not project:
            return False, error_response(
                message=ApiMessage.PROJECT_NOT_FOUND,
                error_code=ErrorCode.RESOURCE_NOT_FOUND,
                status_code=HttpStatus.NOT_FOUND,
            )
        files = project.get("files", [])
        project["files"] = [f for f in files if f.get("file_id") != file_id]
        project["updated_at"] = get_timestamp()
        self._projects[project_id] = project
        self._save_projects()
        return True, None

    def add_chat_history(
        self, project_id: str, message: Dict[str, Any],
    ) -> Tuple[bool, Optional[Tuple]]:
        project = self._projects.get(project_id)
        if not project:
            return False, error_response(
                message=ApiMessage.PROJECT_NOT_FOUND,
                error_code=ErrorCode.RESOURCE_NOT_FOUND,
                status_code=HttpStatus.NOT_FOUND,
            )
        if "timestamp" not in message:
            message["timestamp"] = get_timestamp()
        history = project.setdefault("history", [])
        history.append(message)
        if len(history) > MAX_CHAT_HISTORY_LENGTH:
            project["history"] = history[-MAX_CHAT_HISTORY_LENGTH:]
            logger.debug(f"Trimmed project {project_id} history to {MAX_CHAT_HISTORY_LENGTH} messages")
        project["updated_at"] = get_timestamp()
        self._projects[project_id] = project
        self._save_projects()
        return True, None

    def clear_chat_history(self, project_id: str) -> Tuple[bool, Optional[Tuple]]:
        project = self._projects.get(project_id)
        if not project:
            return False, error_response(
                message=ApiMessage.PROJECT_NOT_FOUND,
                error_code=ErrorCode.RESOURCE_NOT_FOUND,
                status_code=HttpStatus.NOT_FOUND,
            )
        project["history"] = []
        project["updated_at"] = get_timestamp()
        self._projects[project_id] = project
        self._save_projects()
        logger.info(f"Cleared chat history for project: {project_id}")
        return True, None


# ============================================================================
# FILE SERVICE
# ============================================================================

class FileService:
    """
    Manages file uploads, storage, retrieval, and deletion.
    """

    def __init__(self):
        self.upload_dir = UPLOAD_DIRECTORY
        self.temp_dir = TEMP_DIRECTORY
        self.export_dir = EXPORT_DIRECTORY
        ensure_directory(self.upload_dir)
        ensure_directory(self.temp_dir)
        ensure_directory(self.export_dir)
        logger.info(f"File service initialized. Upload dir: {self.upload_dir}")

    def upload_file(
        self, file_data: bytes, filename: str, user_id: str, project_id: Optional[str] = None,
    ) -> Tuple[Optional[Dict[str, Any]], Optional[Tuple]]:
        try:
            if not validate_file_extension(filename):
                return None, error_response(
                    message=ApiMessage.INVALID_FILE_TYPE,
                    error_code=ErrorCode.FILE_UNSUPPORTED_TYPE,
                    status_code=HttpStatus.BAD_REQUEST,
                )
            if not validate_file_size(len(file_data)):
                max_size = humanize_bytes(MAX_UPLOAD_SIZE_BYTES)
                return None, error_response(
                    message=f"{ApiMessage.FILE_TOO_LARGE}. Maximum size: {max_size}",
                    error_code=ErrorCode.FILE_TOO_LARGE,
                    status_code=HttpStatus.BAD_REQUEST,
                )
            safe_filename = sanitize_filename(filename)
            extension = get_file_extension(safe_filename)
            unique_name = f"{generate_id()}{extension}"
            user_dir = self.upload_dir / user_id
            if not ensure_directory(user_dir):
                return None, error_response(
                    message="Failed to create user directory",
                    error_code=ErrorCode.SERVER_INTERNAL,
                    status_code=HttpStatus.INTERNAL_SERVER_ERROR,
                )
            file_path = safe_path_join(user_dir, unique_name)
            if not file_path:
                return None, error_response(
                    message="Invalid file path",
                    error_code=ErrorCode.SERVER_INTERNAL,
                    status_code=HttpStatus.INTERNAL_SERVER_ERROR,
                )
            with open(file_path, "wb") as f:
                f.write(file_data)
            file_info = {
                "file_id": unique_name.rsplit(".", 1)[0] if "." in unique_name else unique_name,
                "original_name": safe_filename,
                "stored_name": unique_name,
                "path": str(file_path),
                "size": len(file_data),
                "size_human": humanize_bytes(len(file_data)),
                "extension": extension,
                "mime_type": get_mime_type(safe_filename),
                "category": get_file_category(safe_filename),
                "user_id": user_id,
                "project_id": project_id,
                "uploaded_at": get_timestamp(),
            }
            logger.info(f"File uploaded: {safe_filename} ({file_info['size_human']})")
            return file_info, None
        except Exception as e:
            logger.error(f"File upload error: {e}", exc_info=True)
            return None, error_response(
                message="File upload failed",
                error_code=ErrorCode.FILE_UPLOAD_FAILED,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def get_file(self, file_path: str) -> Tuple[Optional[bytes], Optional[Tuple]]:
        try:
            path = Path(file_path).resolve()
            if not path.exists():
                return None, error_response(
                    message=ApiMessage.FILE_NOT_FOUND,
                    error_code=ErrorCode.RESOURCE_NOT_FOUND,
                    status_code=HttpStatus.NOT_FOUND,
                )
            upload_dir_resolved = self.upload_dir.resolve()
            if not str(path).startswith(str(upload_dir_resolved)):
                logger.warning(f"Attempted to access file outside upload dir: {path}")
                return None, error_response(
                    message=ApiMessage.FORBIDDEN,
                    error_code=ErrorCode.AUTH_INSUFFICIENT_PERMISSIONS,
                    status_code=HttpStatus.FORBIDDEN,
                )
            with open(path, "rb") as f:
                return f.read(), None
        except Exception as e:
            logger.error(f"File retrieval error: {e}", exc_info=True)
            return None, error_response(
                message="File retrieval failed",
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def delete_file(self, file_path: str) -> Tuple[bool, Optional[Tuple]]:
        try:
            path = Path(file_path).resolve()
            if not path.exists():
                return False, error_response(
                    message=ApiMessage.FILE_NOT_FOUND,
                    error_code=ErrorCode.RESOURCE_NOT_FOUND,
                    status_code=HttpStatus.NOT_FOUND,
                )
            upload_dir_resolved = self.upload_dir.resolve()
            if not str(path).startswith(str(upload_dir_resolved)):
                logger.warning(f"Attempted to delete file outside upload dir: {path}")
                return False, error_response(
                    message=ApiMessage.FORBIDDEN,
                    error_code=ErrorCode.AUTH_INSUFFICIENT_PERMISSIONS,
                    status_code=HttpStatus.FORBIDDEN,
                )
            path.unlink()
            logger.info(f"File deleted: {file_path}")
            return True, None
        except Exception as e:
            logger.error(f"File deletion error: {e}", exc_info=True)
            return False, error_response(
                message="File deletion failed",
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def get_user_files(
        self, user_id: str, category: Optional[str] = None,
    ) -> List[Dict[str, Any]]:
        user_dir = self.upload_dir / user_id
        files = []
        if not user_dir.exists():
            return files
        try:
            for file_path in user_dir.iterdir():
                if not file_path.is_file():
                    continue
                extension = get_file_extension(file_path.name)
                file_category = get_file_category(file_path.name)
                if category and file_category != category:
                    continue
                stat = file_path.stat()
                files.append({
                    "file_id": file_path.stem,
                    "original_name": file_path.name,
                    "stored_name": file_path.name,
                    "path": str(file_path),
                    "size": stat.st_size,
                    "size_human": humanize_bytes(stat.st_size),
                    "extension": extension,
                    "mime_type": get_mime_type(file_path.name),
                    "category": file_category,
                    "user_id": user_id,
                    "uploaded_at": get_utc_now().isoformat(),
                })
        except OSError as e:
            logger.error(f"Error listing user files for {user_id}: {e}")
        files.sort(key=lambda x: x.get("original_name", "").lower())
        return files

    def get_user_storage_stats(self, user_id: str) -> Dict[str, Any]:
        user_dir = self.upload_dir / user_id
        if not user_dir.exists():
            return {"total_files": 0, "total_size_bytes": 0, "total_size_human": "0 B", "categories": {}}
        total_size = 0
        total_files = 0
        categories = {}
        try:
            for file_path in user_dir.iterdir():
                if file_path.is_file():
                    total_files += 1
                    file_size = file_path.stat().st_size
                    total_size += file_size
                    category = get_file_category(file_path.name)
                    if category not in categories:
                        categories[category] = {"count": 0, "size_bytes": 0}
                    categories[category]["count"] += 1
                    categories[category]["size_bytes"] += file_size
        except OSError as e:
            logger.error(f"Error calculating storage stats for {user_id}: {e}")
        return {
            "total_files": total_files,
            "total_size_bytes": total_size,
            "total_size_human": humanize_bytes(total_size),
            "categories": categories,
        }

    def clean_temp_files(self, max_age_hours: int = 24) -> int:
        return clean_temp_files(self.temp_dir, max_age_hours)


# ============================================================================
# EXPORT SERVICE
# ============================================================================

class ExportService:
    """
    Handles exporting content to various formats (PDF, PPTX, TXT, JSON).
    """

    def __init__(self):
        self.export_dir = EXPORT_DIRECTORY
        ensure_directory(self.export_dir)
        logger.info(f"Export service initialized. Export dir: {self.export_dir}")

    def _get_user_export_dir(self, user_id: str) -> Path:
        user_dir = self.export_dir / user_id
        ensure_directory(user_dir)
        return user_dir

    def export_to_pdf(
        self, content: str, filename: str, user_id: str,
    ) -> Tuple[Optional[str], Optional[Tuple]]:
        try:
            user_dir = self._get_user_export_dir(user_id)
            safe_name = sanitize_filename(filename)
            output_path = user_dir / f"{safe_name}.pdf"
            try:
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet
                from reportlab.lib.units import inch
                doc = SimpleDocTemplate(str(output_path), pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
                styles = getSampleStyleSheet()
                story = []
                for line in content.split('\n'):
                    if line.strip():
                        if line.startswith('# '):
                            style, text = styles['Heading1'], line[2:]
                        elif line.startswith('## '):
                            style, text = styles['Heading2'], line[3:]
                        elif line.startswith('### '):
                            style, text = styles['Heading3'], line[4:]
                        else:
                            style, text = styles['Normal'], line
                        story.append(Paragraph(text, style))
                        story.append(Spacer(1, 0.1 * inch))
                doc.build(story)
                logger.info(f"PDF generated with reportlab: {output_path}")
            except ImportError:
                logger.warning("reportlab not installed - saving as text")
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(content)
            return str(output_path), None
        except Exception as e:
            logger.error(f"PDF export error: {e}", exc_info=True)
            return None, error_response(
                message="PDF export failed",
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def export_to_pptx(
        self, slides: List[Dict[str, Any]], filename: str, user_id: str,
    ) -> Tuple[Optional[str], Optional[Tuple]]:
        try:
            user_dir = self._get_user_export_dir(user_id)
            safe_name = sanitize_filename(filename)
            output_path = user_dir / f"{safe_name}.pptx"
            try:
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                for slide_data in slides:
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    if title and slide_data.get("title"):
                        title.text = slide_data["title"]
                    if slide_data.get("content"):
                        slide.placeholders[1].text = slide_data["content"]
                prs.save(str(output_path))
                logger.info(f"PPTX generated: {output_path}")
            except ImportError:
                logger.warning("python-pptx not installed - saving as JSON")
                json_path = output_path.with_suffix(".json")
                with open(json_path, "w", encoding="utf-8") as f:
                    json.dump(slides, f, indent=2, ensure_ascii=False)
                output_path = json_path
            return str(output_path), None
        except Exception as e:
            logger.error(f"PPTX export error: {e}", exc_info=True)
            return None, error_response(
                message="Presentation export failed",
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def export_chat_history(
        self, messages: List[Dict[str, Any]], format: str, filename: str, user_id: str,
    ) -> Tuple[Optional[str], Optional[Tuple]]:
        try:
            user_dir = self._get_user_export_dir(user_id)
            safe_name = sanitize_filename(filename)
            format = format.lower().strip()
            if format == "json":
                output_path = user_dir / f"{safe_name}.json"
                with open(output_path, "w", encoding="utf-8") as f:
                    json.dump(messages, f, indent=2, ensure_ascii=False)
            elif format == "txt":
                output_path = user_dir / f"{safe_name}.txt"
                with open(output_path, "w", encoding="utf-8") as f:
                    for msg in messages:
                        role = msg.get("role", "unknown").upper()
                        content = msg.get("content", "")
                        timestamp = msg.get("timestamp", "")
                        f.write(f"[{role}] {timestamp}\n{content}\n\n")
                        f.write("-" * 60 + "\n\n")
            elif format in ("markdown", "md"):
                output_path = user_dir / f"{safe_name}.md"
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(f"# Chat Export: {safe_name}\n\nExported: {get_timestamp()}\n\n---\n\n")
                    for msg in messages:
                        role = msg.get("role", "unknown").capitalize()
                        content = msg.get("content", "")
                        f.write(f"### {role}\n\n{content}\n\n---\n\n")
            elif format == "pdf":
                text_content = f"Chat Export: {safe_name}\nExported: {get_timestamp()}\n" + "=" * 60 + "\n\n"
                for msg in messages:
                    role = msg.get("role", "unknown").upper()
                    content = msg.get("content", "")
                    text_content += f"[{role}]\n{content}\n\n" + "-" * 40 + "\n\n"
                return self.export_to_pdf(text_content, f"{safe_name}_chat", user_id)
            else:
                return None, error_response(
                    message=f"Unsupported export format: {format}",
                    error_code=ErrorCode.VALIDATION_INVALID_FORMAT,
                    status_code=HttpStatus.BAD_REQUEST,
                )
            logger.info(f"Chat history exported: {output_path}")
            return str(output_path), None
        except Exception as e:
            logger.error(f"Chat export error: {e}", exc_info=True)
            return None, error_response(
                message="Chat export failed",
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def export_document(
        self, content: str, format: str, filename: str, user_id: str,
    ) -> Tuple[Optional[str], Optional[Tuple]]:
        try:
            user_dir = self._get_user_export_dir(user_id)
            safe_name = sanitize_filename(filename)
            format = format.lower().strip()
            if format == "pdf":
                return self.export_to_pdf(content, safe_name, user_id)
            elif format == "txt":
                output_path = user_dir / f"{safe_name}.txt"
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(content)
                return str(output_path), None
            elif format == "docx":
                output_path = user_dir / f"{safe_name}.docx"
                try:
                    from docx import Document
                    doc = Document()
                    for line in content.split('\n'):
                        doc.add_paragraph(line)
                    doc.save(str(output_path))
                    logger.info(f"DOCX generated: {output_path}")
                except ImportError:
                    logger.warning("python-docx not installed - saving as text")
                    output_path = output_path.with_suffix(".txt")
                    with open(output_path, "w", encoding="utf-8") as f:
                        f.write(content)
                return str(output_path), None
            elif format in ("markdown", "md"):
                output_path = user_dir / f"{safe_name}.md"
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(content)
                return str(output_path), None
            else:
                return None, error_response(
                    message=f"Unsupported export format: {format}",
                    error_code=ErrorCode.VALIDATION_INVALID_FORMAT,
                    status_code=HttpStatus.BAD_REQUEST,
                )
        except Exception as e:
            logger.error(f"Document export error: {e}", exc_info=True)
            return None, error_response(
                message="Document export failed",
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )


# ============================================================================
# PDF ANALYSIS SERVICE
# ============================================================================

class PDFAnalysisService:
    """
    Handles PDF document analysis and text extraction.
    """

    def __init__(self):
        self._pypdf = _import_pypdf()
        self._pdfplumber = _import_pdfplumber()
        if self._pypdf:
            logger.info("PDF extraction: PyPDF2 available")
        if self._pdfplumber:
            logger.info("PDF extraction: pdfplumber available")
        if not self._pypdf and not self._pdfplumber:
            logger.warning("No PDF extraction library available.")

    def extract_text(self, file_path: str) -> Tuple[Optional[str], Optional[Tuple]]:
        try:
            path = Path(file_path).resolve()
            if not path.exists():
                return None, error_response(
                    message=ApiMessage.FILE_NOT_FOUND,
                    error_code=ErrorCode.RESOURCE_NOT_FOUND,
                    status_code=HttpStatus.NOT_FOUND,
                )
            if path.suffix.lower() != '.pdf':
                return None, error_response(
                    message="File is not a PDF",
                    error_code=ErrorCode.FILE_UNSUPPORTED_TYPE,
                    status_code=HttpStatus.BAD_REQUEST,
                )
            file_size_mb = path.stat().st_size / (1024 * 1024)
            if file_size_mb > 50:
                return None, error_response(
                    message=f"PDF file too large ({file_size_mb:.1f}MB). Maximum: 50MB",
                    error_code=ErrorCode.FILE_TOO_LARGE,
                    status_code=HttpStatus.BAD_REQUEST,
                )
            extracted_text = ""
            if self._pdfplumber:
                try:
                    with self._pdfplumber.open(path) as pdf:
                        total_pages = len(pdf.pages)
                        if total_pages > SystemLimits.MAX_PDF_PAGES:
                            return None, error_response(
                                message=f"PDF has too many pages ({total_pages}).",
                                error_code=ErrorCode.FILE_UNSUPPORTED_TYPE,
                                status_code=HttpStatus.BAD_REQUEST,
                            )
                        for page_num, page in enumerate(pdf.pages, 1):
                            try:
                                page_text = page.extract_text()
                                if page_text:
                                    extracted_text += page_text + "\n"
                            except Exception as e:
                                logger.warning(f"Error extracting text from page {page_num}: {e}")
                        if extracted_text.strip():
                            logger.info(f"PDF text extracted with pdfplumber: {total_pages} pages, {len(extracted_text)} chars")
                            return extracted_text.strip(), None
                except Exception as e:
                    logger.warning(f"pdfplumber extraction failed: {e}")
            if self._pypdf:
                try:
                    with open(path, "rb") as f:
                        reader = self._pypdf.PdfReader(f)
                        total_pages = len(reader.pages)
                        if total_pages > SystemLimits.MAX_PDF_PAGES:
                            return None, error_response(
                                message=f"PDF has too many pages ({total_pages}).",
                                error_code=ErrorCode.FILE_UNSUPPORTED_TYPE,
                                status_code=HttpStatus.BAD_REQUEST,
                            )
                        if reader.is_encrypted:
                            try:
                                reader.decrypt("")
                            except Exception:
                                return None, error_response(
                                    message="PDF is encrypted and cannot be read",
                                    error_code=ErrorCode.FILE_CORRUPTED,
                                    status_code=HttpStatus.BAD_REQUEST,
                                )
                        for page_num, page in enumerate(reader.pages, 1):
                            try:
                                page_text = page.extract_text()
                                if page_text:
                                    extracted_text += page_text + "\n"
                            except Exception as e:
                                logger.warning(f"Error extracting text from page {page_num}: {e}")
                        if extracted_text.strip():
                            logger.info(f"PDF text extracted with PyPDF2: {total_pages} pages, {len(extracted_text)} chars")
                            return extracted_text.strip(), None
                except Exception as e:
                    logger.warning(f"PyPDF2 extraction failed: {e}")
            if not extracted_text.strip():
                return None, error_response(
                    message="Could not extract text from PDF. The file may be scanned/image-based.",
                    error_code=ErrorCode.FILE_CORRUPTED,
                    status_code=HttpStatus.BAD_REQUEST,
                )
            return extracted_text.strip(), None
        except Exception as e:
            logger.error(f"PDF extraction error: {e}", exc_info=True)
            return None, error_response(
                message="Failed to extract text from PDF",
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def extract_metadata(self, file_path: str) -> Tuple[Optional[Dict[str, Any]], Optional[Tuple]]:
        try:
            path = Path(file_path).resolve()
            if not path.exists():
                return None, error_response(
                    message=ApiMessage.FILE_NOT_FOUND,
                    error_code=ErrorCode.RESOURCE_NOT_FOUND,
                    status_code=HttpStatus.NOT_FOUND,
                )
            metadata = {
                "filename": path.name,
                "size_bytes": path.stat().st_size,
                "size_human": humanize_bytes(path.stat().st_size),
                "extension": path.suffix.lower(),
                "mime_type": get_mime_type(path.name),
            }
            if self._pypdf:
                try:
                    with open(path, "rb") as f:
                        reader = self._pypdf.PdfReader(f)
                        metadata["total_pages"] = len(reader.pages)
                        metadata["is_encrypted"] = reader.is_encrypted
                        if reader.metadata:
                            pdf_meta = reader.metadata
                            for key in ("title", "author", "subject", "creator", "producer"):
                                metadata[key] = getattr(pdf_meta, key, None)
                except Exception as e:
                    logger.debug(f"Could not extract PDF metadata: {e}")
            elif self._pdfplumber:
                try:
                    with self._pdfplumber.open(path) as pdf:
                        metadata["total_pages"] = len(pdf.pages)
                        if pdf.metadata:
                            metadata.update({k: v for k, v in pdf.metadata.items() if v})
                except Exception as e:
                    logger.debug(f"Could not extract PDF metadata with pdfplumber: {e}")
            return metadata, None
        except Exception as e:
            logger.error(f"PDF metadata extraction error: {e}", exc_info=True)
            return None, error_response(
                message="Failed to extract PDF metadata",
                error_code=ErrorCode.SERVER_INTERNAL,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )

    def analyze_with_ai(
        self, file_path: str, ai_service: 'AIService', analysis_type: str = "summarize", model: str = None,
    ) -> Tuple[Optional[Dict[str, Any]], Optional[Tuple]]:
        try:
            if model is None:
                model = AIModel.GITHUB_GPT4O.value
            text, error = self.extract_text(file_path)
            if error:
                return None, error
            if not text:
                return None, error_response(
                    message="No text could be extracted from the PDF",
                    error_code=ErrorCode.FILE_CORRUPTED,
                    status_code=HttpStatus.BAD_REQUEST,
                )
            max_chars = 30000
            truncated = False
            if len(text) > max_chars:
                text = text[:max_chars] + "\n\n[Document truncated due to length...]"
                truncated = True
            analysis_prompts = {
                "summarize": "Please provide a comprehensive yet concise summary of the following document. Include the main topics, key findings, and important conclusions.\n\nDocument:\n\n",
                "extract": "Extract key information, facts, figures, dates, and names from this document. Present the information in a structured, easy-to-read format.\n\nDocument:\n\n",
                "analyze": "Analyze this document thoroughly. Identify the main themes, arguments, supporting evidence, potential biases, and overall effectiveness. Provide critical insights.\n\nDocument:\n\n",
                "keywords": "Extract the most important keywords, topics, and concepts from this document. Group related terms together and indicate their relevance.\n\nDocument:\n\n",
                "questions": "Based on this document, generate a list of insightful questions that would help deepen understanding of the content. Include both factual and analytical questions.\n\nDocument:\n\n",
            }
            prompt = analysis_prompts.get(analysis_type, analysis_prompts["summarize"]) + text
            messages = [{"role": "user", "content": prompt}]
            result, error = ai_service.chat(
                messages=messages, model=model,
                system_prompt=DefaultPrompts.PDF_ANALYZER, max_tokens=3000, temperature=0.3,
            )
            if error:
                return None, error
            metadata, _ = self.extract_metadata(file_path)
            analysis_result = {
                "analysis_type": analysis_type,
                "content": result.get("content", ""),
                "model": result.get("model"),
                "provider": result.get("provider"),
                "usage": result.get("usage", {}),
                "text_length": len(text),
                "text_truncated": truncated,
                "document_metadata": metadata,
            }
            logger.info(f"PDF analysis completed: {analysis_type} using {model}")
            return analysis_result, None
        except Exception as e:
            logger.error(f"PDF analysis error: {e}", exc_info=True)
            return None, error_response(
                message="PDF analysis failed",
                error_code=ErrorCode.AI_GENERATION_ERROR,
                status_code=HttpStatus.INTERNAL_SERVER_ERROR,
            )


# ============================================================================
# GLOBAL SERVICE INSTANCES
# ============================================================================

_ai_service_instance = None
_project_service_instance = None
_file_service_instance = None
_export_service_instance = None
_pdf_service_instance = None


def get_ai_service() -> AIService:
    global _ai_service_instance
    if _ai_service_instance is None:
        _ai_service_instance = AIService()
    return _ai_service_instance


def get_project_service() -> ProjectService:
    global _project_service_instance
    if _project_service_instance is None:
        _project_service_instance = ProjectService()
    return _project_service_instance


def get_file_service() -> FileService:
    global _file_service_instance
    if _file_service_instance is None:
        _file_service_instance = FileService()
    return _file_service_instance


def get_export_service() -> ExportService:
    global _export_service_instance
    if _export_service_instance is None:
        _export_service_instance = ExportService()
    return _export_service_instance


def get_pdf_service() -> PDFAnalysisService:
    global _pdf_service_instance
    if _pdf_service_instance is None:
        _pdf_service_instance = PDFAnalysisService()
    return _pdf_service_instance


ai_service = get_ai_service()
project_service = get_project_service()
file_service = get_file_service()
export_service = get_export_service()
pdf_service = get_pdf_service()